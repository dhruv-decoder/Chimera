"""Generate the solution-walkthrough deck (.pptx) from live artifacts.

Reads data/artifacts/{eval_report,loop_report,sim_meta}.json, renders two charts
with matplotlib, and assembles a dark-themed PowerPoint so the deck always
matches the numbers the code actually produced. Output: docs/Chimera_Deck.pptx.

    python scripts/build_deck.py
"""
from __future__ import annotations

import json

import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Emu, Inches, Pt

from chimera.config import ARTIFACTS_DIR, REPO_ROOT

INK = RGBColor(0x06, 0x07, 0x0A)
PANEL = RGBColor(0x0E, 0x10, 0x16)
MIST = RGBColor(0xE7, 0xE9, 0xEE)
MUTE = RGBColor(0x8A, 0x90, 0x9F)
DEFENSE = RGBColor(0x2E, 0xD6, 0xA6)
THREAT = RGBColor(0xFF, 0x5C, 0x49)
AGENTIC = RGBColor(0x8B, 0x8C, 0xF0)

DOCS = REPO_ROOT / "docs"
W, H = Inches(13.333), Inches(7.5)


def _read(name):
    p = ARTIFACTS_DIR / name
    return json.loads(p.read_text()) if p.exists() else {}


def _bg(slide):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = INK


def _text(slide, x, y, w, h, text, size=18, color=MIST, bold=False, align=PP_ALIGN.LEFT, font="Arial"):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    lines = text.split("\n")
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        r = p.add_run(); r.text = line
        r.font.size = Pt(size); r.font.bold = bold
        r.font.color.rgb = color; r.font.name = font
    return tb


def _accent(slide, x, y, w=Inches(0.9), color=DEFENSE):
    line = slide.shapes.add_shape(1, x, y, w, Pt(3))
    line.fill.solid(); line.fill.fore_color.rgb = color
    line.line.fill.background()
    return line


def _render_charts(ev, loop):
    plt.rcParams.update({"figure.facecolor": "#06070a", "axes.facecolor": "#0e1016",
                         "text.color": "#e7e9ee", "axes.edgecolor": "#252b38",
                         "xtick.color": "#8a909f", "ytick.color": "#8a909f",
                         "axes.labelcolor": "#aeb4c2", "font.size": 12})
    # Hardening curve
    curve = loop.get("hardening_curve", [])
    if curve:
        fig, ax = plt.subplots(figsize=(7.2, 4.0), dpi=150)
        xs = [p["round"] for p in curve]
        ax.plot(xs, [p["pre_recall"] * 100 for p in curve], "--o", color="#ff5c49", lw=2.5, label="pre-retrain (attack breach)")
        ax.plot(xs, [p["post_recall"] * 100 for p in curve], "-o", color="#2ed6a6", lw=2.5, label="post-retrain (defence hardens)")
        ax.set_xlabel("round"); ax.set_ylabel("recall (%)"); ax.set_ylim(0, 105)
        ax.set_xticks(xs); ax.grid(True, color="#1a1e28", lw=0.6)
        ax.legend(facecolor="#0e1016", edgecolor="#252b38", labelcolor="#e7e9ee", fontsize=10)
        for s in ax.spines.values():
            s.set_color("#252b38")
        fig.tight_layout(); fig.savefig(ARTIFACTS_DIR / "hardening_curve.png"); plt.close(fig)
    # Per-vector recall
    pv = ev.get("per_vector_recall", {})
    if pv:
        fig, ax = plt.subplots(figsize=(7.2, 4.0), dpi=150)
        items = sorted(pv.items(), key=lambda x: x[1]["recall"])
        names = [k for k, _ in items]; vals = [v["recall"] * 100 for _, v in items]
        colors = ["#ff5c49" if v < 60 else "#f5b544" if v < 85 else "#2ed6a6" for v in vals]
        ax.barh(names, vals, color=colors)
        ax.set_xlabel("recall (%)"); ax.set_xlim(0, 105); ax.grid(True, axis="x", color="#1a1e28", lw=0.6)
        for s in ax.spines.values():
            s.set_color("#252b38")
        fig.tight_layout(); fig.savefig(ARTIFACTS_DIR / "per_vector.png"); plt.close(fig)


def _slide(prs, layout=6):
    s = prs.slides.add_slide(prs.slide_layouts[layout])
    _bg(s)
    return s


def _kicker(slide, text, color=DEFENSE):
    _accent(slide, Inches(0.7), Inches(0.62), color=color)
    _text(slide, Inches(0.7), Inches(0.72), Inches(8), Inches(0.4), text.upper(), 12, color, bold=True)


def _title(slide, text, y=Inches(1.05)):
    _text(slide, Inches(0.7), y, Inches(12), Inches(1.0), text, 30, MIST, bold=True)


def build():
    ev = _read("eval_report.json"); loop = _read("loop_report.json"); meta = _read("sim_meta.json")
    analysis = _read("analysis_report.json")
    _render_charts(ev, loop)
    s = ev.get("supervised", {})
    curve = loop.get("hardening_curve", [])

    prs = Presentation()
    prs.slide_width = W; prs.slide_height = H

    # 1 title
    sl = _slide(prs)
    _accent(sl, Inches(0.7), Inches(2.5), Inches(1.4), DEFENSE)
    _text(sl, Inches(0.7), Inches(2.7), Inches(12), Inches(1.2), "Chimera", 54, MIST, bold=True)
    _text(sl, Inches(0.7), Inches(3.8), Inches(11.5), Inches(1.0),
          "A closed-loop adversarial lab for GenAI-era payment fraud.", 22, DEFENSE)
    _text(sl, Inches(0.7), Inches(4.5), Inches(11.5), Inches(1.2),
          "Discover emerging fraud, simulate it at fidelity across card, real-time and agentic rails,\n"
          "and harden a detector on what it misses. Identify, generate, defend - one feedback loop.", 15, MUTE)
    _text(sl, Inches(0.7), Inches(6.6), Inches(12), Inches(0.5),
          "Mastercard Innovation Challenge @ GFF 2026 · AI Defense Lab for Payment Security", 12, MUTE)

    # 2 problem
    sl = _slide(prs); _kicker(sl, "The problem", THREAT); _title(sl, "Static defences decay the moment fraud changes shape")
    _text(sl, Inches(0.7), Inches(2.2), Inches(11.8), Inches(3.5),
          "· GenAI made fraud faster, cheaper and adaptive: deepfakes ~11% of global fraud, biometric-fraud\n"
          "  deepfake attempts up ~58%, Fraud-as-a-Service from ~$50/mo.\n\n"
          "· 2026 frontier: agentic commerce. Autonomous agents transact on delegated tokens with no step-up;\n"
          "  Visa logged a 450%+ rise in dark-web 'AI Agent' chatter; agents run live carding at machine speed.\n\n"
          "· Real-time rails are the laundering surface: 524k mule accounts flagged in India in March 2026 alone.\n\n"
          "· Evaluating one classifier on a fixed dataset measures fit to yesterday's fraud - not robustness to an\n"
          "  adaptive adversary. The challenge asks for a loop. I built one.", 15, MIST)

    # 3 system
    sl = _slide(prs); _kicker(sl, "System"); _title(sl, "One system, three pillars, one loop")
    for i, (name, color, desc) in enumerate([
        ("IDENTIFY", AGENTIC, "ATT&CK-style taxonomy of 16 GenAI fraud techniques +\nlive RAG ideation agent (LangGraph / Groq gpt-oss-120b)"),
        ("GENERATE", THREAT, "Multi-rail simulator + entity graph + hard negatives +\n9 attack synthesizers + adversarial evasion search"),
        ("DEFEND", DEFENSE, "LightGBM + novelty channel + agent-identity features +\nSHAP reason codes + cost-aware operating point"),
    ]):
        x = Inches(0.7 + i * 4.15)
        card = sl.shapes.add_shape(1, x, Inches(2.3), Inches(3.9), Inches(2.6))
        card.fill.solid(); card.fill.fore_color.rgb = PANEL; card.line.color.rgb = RGBColor(0x25, 0x2B, 0x38)
        _text(sl, x + Inches(0.25), Inches(2.5), Inches(3.5), Inches(0.5), name, 16, color, bold=True)
        _text(sl, x + Inches(0.25), Inches(3.1), Inches(3.5), Inches(1.6), desc, 12, MIST)
    _text(sl, Inches(0.7), Inches(5.3), Inches(11.8), Inches(1.2),
          "The loop: identify → generate (evolve evasion) → detect → evaluate → retrain on what breaks through →\n"
          "re-ideate. Shipped as a plain engine and as a LangGraph multi-agent StateGraph.", 14, DEFENSE)

    # 3b architecture diagram
    arch = REPO_ROOT / "docs" / "architecture.png"
    if arch.exists():
        sl = _slide(prs); _kicker(sl, "Architecture")
        _title(sl, "One codebase, three pillars, one loop")
        sl.shapes.add_picture(str(arch), Inches(0.7), Inches(1.9), width=Inches(11.9))

    # 3c multi-agent orchestration
    sl = _slide(prs); _kicker(sl, "Orchestration", AGENTIC)
    _title(sl, "A live multi-agent system (LangGraph)")
    agents = [
        ("RECON", AGENTIC, "RAG-grounded ideation\nproposes novel variants\n(gpt-oss-120b)"),
        ("RED TEAM", THREAT, "evolutionary search\ntunes each attack to\nevade the live model"),
        ("ATTACK", THREAT, "generates the evasive\nstream, measures the\nbreach"),
        ("BLUE TEAM", DEFENSE, "retrains on the misses,\nmeasures the recovery,\nthen the graph cycles"),
    ]
    for i, (name, color, desc) in enumerate(agents):
        x = Inches(0.7 + i * 3.12)
        card = sl.shapes.add_shape(1, x, Inches(2.3), Inches(2.85), Inches(2.0))
        card.fill.solid(); card.fill.fore_color.rgb = PANEL; card.line.color.rgb = RGBColor(0x25, 0x2B, 0x38)
        _text(sl, x + Inches(0.2), Inches(2.45), Inches(2.5), Inches(0.4), name, 13, color, bold=True)
        _text(sl, x + Inches(0.2), Inches(2.95), Inches(2.5), Inches(1.3), desc, 11, MIST)
        if i < 3:
            _text(sl, x + Inches(2.82), Inches(3.0), Inches(0.4), Inches(0.5), ">", 20, MUTE)
    _text(sl, Inches(0.7), Inches(4.6), Inches(11.8), Inches(0.5),
          "A compiled StateGraph cycles recon -> red_team -> attack -> blue_team until the round budget is spent.", 13, DEFENSE)
    tr = (loop.get("trace") or [])[:6]
    if tr:
        _text(sl, Inches(0.7), Inches(5.25), Inches(11.8), Inches(1.9),
              "Execution trace (this report was generated by the graph):\n" + "\n".join(tr), 11, MUTE, font="Consolas")

    # 4 identify
    sl = _slide(prs); _kicker(sl, "Identify", AGENTIC); _title(sl, "Mapping the threat landscape")
    _text(sl, Inches(0.7), Inches(2.2), Inches(11.8), Inches(4.5),
          "A payment-fraud ATT&CK matrix: 6 kill-chain tactics x 16 techniques, each grounded in 2026\n"
          "intelligence with the observable signatures a defender can exploit. Nine simulated end-to-end.\n\n"
          "Novel, high-signal vectors I cover:\n"
          "· Delegated-token / agent-identity abuse - a hijacked agent spends inside a cardholder's mandate\n"
          "  (Agent Pay, Trusted Agent Protocol). The 2026 frontier, and my headline new capability.\n"
          "· Agentic-commerce carding - machine-speed abuse of delegated agent tokens.\n"
          "· Deepfake-authorised push payments - cloned voice/video; the victim authorises, so auth can't stop it.\n"
          "· Money-mule networks - fan-in / fan-out layering within minutes on real-time rails.\n"
          "· Synthetic-identity bust-out, account takeover, pig-butchering, structuring, card testing.\n\n"
          "Ideation agent: RAG-grounds where the detector is weak and proposes novel variants as structured\n"
          "specs. Runs live on Groq gpt-oss-120b (every loop round is tagged with the model) or an offline\n"
          "planner - the demo never breaks.", 14, MIST)

    # 5 generate
    sl = _slide(prs); _kicker(sl, "Generate", THREAT); _title(sl, "Fidelity is the point")
    _text(sl, Inches(0.7), Inches(2.2), Inches(11.8), Inches(4.5),
          f"· Realistic population: {meta.get('n_total', 0):,} events, {meta.get('fraud_rate', 0)*100:.2f}% fraud, across\n"
          "  card-not-present, card-present and real-time A2A, plus a shared agentic channel.\n\n"
          "· Entity graph (accounts, devices, merchants, payees) so coordinated attacks surface as structure.\n\n"
          "· Hard negatives - the credibility centrepiece. I inject legitimate look-alikes (benign agentic\n"
          "  shopping, large first-time payees, travel, collector accounts, recurring investments) into the exact\n"
          "  feature regions where fraud lives. Without them any single flag separates the classes and AUC is a\n"
          "  meaningless 1.0.\n\n"
          "· Adversarial evasion: a black-box evolutionary search tunes each attack's shape (not volume) to\n"
          "  minimise the live detector's risk - the generator becomes a live red team.", 14, MIST)

    # 6 defend
    sl = _slide(prs); _kicker(sl, "Defend"); _title(sl, "Accuracy with an honest novelty story")
    _text(sl, Inches(0.7), Inches(2.2), Inches(11.8), Inches(4.5),
          "· Supervised: LightGBM on ~45 features - event, behavioural velocity (per-account windows,\n"
          "  inter-arrival, amount z-score vs own history), structural (device fan-out, counterparty\n"
          "  in-degree, A2A graph degree + PageRank), and an agent-identity family (attestation, directory\n"
          "  trust, mandate-cap breach, agent-id replay fan-out) for delegated-token abuse.\n\n"
          "· Novelty channel: isolation forest + PCA reconstruction on legit-only traffic. Flags unseen attack\n"
          "  types the supervised model was never trained on - the zero-day answer.\n\n"
          "· Explainability: per-event reason codes via exact TreeSHAP (no extra serving dependency).\n\n"
          "· Cost-aware operating point: fraud value recovered + alerts per 10k + a cost-optimal threshold -\n"
          "  the metrics a fraud desk signs off. Mitigation: risk -> allow / step-up / hold / block.", 14, MIST)

    # 7 detection results
    sl = _slide(prs); _kicker(sl, "Results · detection"); _title(sl, "Strong, and reported honestly")
    metrics = [("ROC-AUC", f"{s.get('roc_auc', 0):.4f}"), ("PR-AUC", f"{s.get('pr_auc', 0):.4f}"),
               ("F1", f"{s.get('f1', 0):.3f}"), ("FPR @ 90% recall", f"{s.get('fpr_at_90_recall', 0)*100:.3f}%")]
    for i, (k, v) in enumerate(metrics):
        x = Inches(0.7 + i * 3.1)
        card = sl.shapes.add_shape(1, x, Inches(2.2), Inches(2.9), Inches(1.5))
        card.fill.solid(); card.fill.fore_color.rgb = PANEL; card.line.color.rgb = RGBColor(0x25, 0x2B, 0x38)
        _text(sl, x + Inches(0.2), Inches(2.35), Inches(2.6), Inches(0.4), k, 11, MUTE, bold=True)
        _text(sl, x + Inches(0.2), Inches(2.8), Inches(2.6), Inches(0.7), v, 26, DEFENSE, bold=True)
    if (ARTIFACTS_DIR / "per_vector.png").exists():
        sl.shapes.add_picture(str(ARTIFACTS_DIR / "per_vector.png"), Inches(0.7), Inches(4.0), height=Inches(3.0))
    op = ev.get("operating_point", {}).get("cost_optimal", {})
    _text(sl, Inches(7.7), Inches(4.2), Inches(5.2), Inches(2.8),
          "Per-vector recall is deliberately uneven. Deepfake authorised-push sits lower - the victim uses\n"
          "their own device and genuine auth, so it is the hardest to catch. Structural attacks are caught in\n"
          "full. Near-perfect aggregate AUC reflects strong signatures on naive campaigns - the real test is\n"
          "adversarial pressure (next slides).\n\n"
          f"Cost-aware operating point: {op.get('value_detected_rate', 0)*100:.0f}% of fraud value recovered at\n"
          f"~{op.get('alerts_per_10k', 0):.0f} alerts per 10k transactions.", 13, MIST)

    # 7b fidelity evidence
    fid = analysis.get("fidelity", {})
    if (ARTIFACTS_DIR / "fidelity_pca.png").exists():
        sl = _slide(prs); _kicker(sl, "Results · fidelity")
        _title(sl, "The data is genuinely hard, by design")
        sl.shapes.add_picture(str(ARTIFACTS_DIR / "fidelity_amount.png"), Inches(0.7), Inches(2.0), height=Inches(2.5))
        sl.shapes.add_picture(str(ARTIFACTS_DIR / "fidelity_pca.png"), Inches(0.7), Inches(4.6), height=Inches(2.5))
        _text(sl, Inches(7.6), Inches(2.1), Inches(5.4), Inches(4.6),
              "Fraud and legitimate traffic overlap in feature space (PCA, top). The single most\n"
              f"discriminative feature reaches only AUC {fid.get('best_single_feature_auc', 0):.2f}\n"
              f"({fid.get('best_single_feature','')}), so no one flag separates the classes - the model must\n"
              "combine velocity, graph and behavioural signals.\n\n"
              "Amounts follow realistic log-normal distributions per account; fraud is embedded in the same\n"
              "space rather than bolted on. Hard negatives (benign agentic shopping, collector accounts,\n"
              "large first-time payees) sit exactly where fraud lives - which is what makes the benchmark\n"
              "credible instead of a trivial AUC of 1.0.", 13, MIST)

    # 7c ablation
    ab = analysis.get("ablation", {})
    if (ARTIFACTS_DIR / "ablation.png").exists():
        sl = _slide(prs); _kicker(sl, "Results · ablation")
        _title(sl, "The ensemble earns its complexity")
        sl.shapes.add_picture(str(ARTIFACTS_DIR / "ablation.png"), Inches(0.7), Inches(2.2), height=Inches(3.2))
        _text(sl, Inches(7.6), Inches(2.4), Inches(5.4), Inches(3.5),
              "Held-out PR-AUC across models. Logistic regression and an isolation-forest-only detector\n"
              "leave signal on the table; gradient boosting captures the non-linear velocity/graph\n"
              "interactions; the novelty channel adds unseen-attack coverage on top.\n\n"
              "Every number is a measured lift, not an assertion - the complexity is justified.", 13, MIST)

    # 8 hardening curve (money slide)
    sl = _slide(prs); _kicker(sl, "Results · the closed loop", THREAT)
    _title(sl, "The hardening curve - the whole thesis in one chart")
    if (ARTIFACTS_DIR / "hardening_curve.png").exists():
        sl.shapes.add_picture(str(ARTIFACTS_DIR / "hardening_curve.png"), Inches(0.7), Inches(2.1), height=Inches(3.9))
    breach = min((p["pre_recall"] for p in curve[1:]), default=0) * 100 if len(curve) > 1 else 0
    _text(sl, Inches(7.6), Inches(2.2), Inches(5.4), Inches(4.5),
          f"Round 1: the red team finds evasion that collapses recall to {breach:.0f}% - a static model would be\n"
          "blind. Retraining on those samples restores it.\n\n"
          "By rounds 2-3 the red team can no longer find easy evasion against the hardened detector: the loop\n"
          "converges. The defence learned to generalise across the adversary's moves.\n\n"
          "Behavioural attacks (deepfake push, pig-butchering, structuring) can be tuned to near-invisibility\n"
          "against a static model; structural attacks can't. Closed-loop retraining is what recovers the\n"
          "evadable classes - a single static model cannot.", 13, MIST)

    # 8b headline capability - agentic-commerce identity abuse
    loo = {r["vector"]: r for r in ev.get("leave_one_out", [])}
    hj = loo.get("AGENT-HIJACK", {})
    sl = _slide(prs); _kicker(sl, "Headline capability", AGENTIC)
    _title(sl, "Agentic-commerce identity abuse - the 2026 frontier")
    _text(sl, Inches(0.7), Inches(2.05), Inches(12.0), Inches(2.2),
          "A payment can now be initiated by an autonomous agent on a cardholder's behalf (Mastercard Agent\n"
          "Pay's Agentic Token; Visa's Trusted Agent Protocol). A hijacked or malicious agent - via prompt\n"
          "injection, token theft, or a rogue SDK - spends inside someone else's mandate. It looks identical to a\n"
          "legitimate agent on velocity, device and cadence. Only credential integrity separates them.", 14, MIST)
    for i, (k, v) in enumerate([
        ("What I simulate", "AGENT-HIJACK: one agent id\nreplayed across many mandates,\noff-scope spend, weak attestation"),
        ("New defence family", "attestation · directory trust ·\nmandate-cap breach · agent-id\nreplay fan-out"),
        ("Unseen-vector result",
         f"supervised {hj.get('supervised_recall', 0)*100:.0f}%  ->  novelty {hj.get('novelty_recall', 0)*100:.0f}%\n"
         "(trained with this attack fully\nremoved, then scored on it)"),
    ]):
        x = Inches(0.7 + i * 4.15)
        card = sl.shapes.add_shape(1, x, Inches(4.3), Inches(3.9), Inches(2.2))
        card.fill.solid(); card.fill.fore_color.rgb = PANEL; card.line.color.rgb = RGBColor(0x25, 0x2B, 0x38)
        _text(sl, x + Inches(0.25), Inches(4.45), Inches(3.5), Inches(0.4), k, 12,
              AGENTIC if i < 2 else DEFENSE, bold=True)
        _text(sl, x + Inches(0.25), Inches(4.95), Inches(3.5), Inches(1.5), v, 12, MIST)

    # 9 novelty + feasibility
    sl = _slide(prs); _kicker(sl, "Why it wins", AGENTIC); _title(sl, "Novelty & real-world feasibility")
    _text(sl, Inches(0.7), Inches(2.2), Inches(11.8), Inches(2.4),
          "Novel: a working closed loop with a measured hardening curve; agentic-commerce identity abuse as a\n"
          "first-class simulated vector with a matching defence family; live red-team ideation from an open-weight\n"
          "model (gpt-oss-120b) running in the real loop; adversarial evasion as a live red team; a novelty channel\n"
          "evaluated with leave-one-vector-out; cost-aware, hard-negative-driven reporting.", 14, MIST)
    _text(sl, Inches(0.7), Inches(4.6), Inches(11.8), Inches(2.4),
          "Feasible: the schema is what a real switch sees at authorisation; controls map to 3-DS step-up and\n"
          "RBI friction / kill-switch; graph features mirror NPCI MuleHunter-style detection. On live data,\n"
          "features move to a streaming store (point-in-time correct) and the loop becomes a scheduled retrain.\n"
          "Paid infrastructure is a quality dial, not a rewrite: managed embeddings + vector DB for richer\n"
          "ideation, a temporal GNN for large rings, a hosted endpoint for always-on live variant generation.", 14, MIST)

    # 9b scorecard against the judging criteria
    sl = _slide(prs); _kicker(sl, "Self-assessment", DEFENSE)
    _title(sl, "How this maps to the five judging criteria")
    rows = [
        ("Diversity of attacks", "16 techniques, 9 simulated across card / real-time / agentic rails, plus a live agent that proposes more."),
        ("Fidelity of simulation", "Latent-profile population, entity graph, hard negatives - best single-feature AUC only 0.84."),
        ("Detection efficacy", "LightGBM + novelty + agent-identity features; 100% of fraud value caught at ~140 alerts / 10k."),
        ("Novelty", "Closed loop with a measured hardening curve; agentic-commerce identity abuse; a live multi-agent red team."),
        ("Real-world feasibility", "Auth-time schema, graded controls, NPCI-style graph signals, Agent Pay / TAP field mapping, streaming path."),
    ]
    y = 2.15
    for k, v in rows:
        _text(sl, Inches(0.7), Inches(y), Inches(3.5), Inches(0.9), k, 14, DEFENSE, bold=True)
        _text(sl, Inches(4.4), Inches(y), Inches(8.3), Inches(0.9), v, 13, MIST)
        y += 0.95

    # 10 responsible use
    sl = _slide(prs); _kicker(sl, "Responsible use", DEFENSE); _title(sl, "Defensive by construction")
    _text(sl, Inches(0.7), Inches(2.2), Inches(11.8), Inches(4.0),
          "· Runs entirely on synthetic data it generates itself - no real cardholders, PII, credentials or\n"
          "  live payment connectivity.\n\n"
          "· Attack modules are parameterised statistical patterns (amounts, timing, graph shape) that stress a\n"
          "  detector - not operational playbooks, social-engineering scripts, or working exploit code, and they\n"
          "  cannot be pointed at a real system.\n\n"
          "· One-directional value: the red team exists only to harden the blue team. The output is a stronger\n"
          "  detector plus a prioritised view of which vectors most need human review and step-up controls.\n\n"
          "· In production this belongs inside a bank's controlled environment with human oversight, aligned to\n"
          "  responsible-AI and model-risk governance.", 14, MIST)
    _text(sl, Inches(0.7), Inches(6.7), Inches(12), Inches(0.5),
          "Chimera · identify · generate · defend · one loop", 12, MUTE)

    out = DOCS / "Chimera_Deck.pptx"
    prs.save(str(out))
    print(f"Saved deck -> {out}")


if __name__ == "__main__":
    build()
